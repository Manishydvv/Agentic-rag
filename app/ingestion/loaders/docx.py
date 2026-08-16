"""
Office Document Loaders
=======================
Supports all major Microsoft Office formats:
  - .docx  → Word documents       (python-docx)
  - .xlsx  → Excel spreadsheets   (openpyxl)
  - .pptx  → PowerPoint files     (python-pptx)
  - .csv   → CSV data files       (built-in csv)

All parsers return List[Document] with source metadata.
"""

import csv
import os
from typing import List
from langchain_core.documents import Document
from app.utils.logger import logger


# ──────────────────────────────────────────────────────────────────────────────
# DOCX — Word Documents
# ──────────────────────────────────────────────────────────────────────────────
def parse_docx(file_path: str) -> List[Document]:
    """
    Parses a Word .docx file.
    Returns one Document per section (split by headings).
    Falls back to one Document for the entire file if no headings found.
    """
    try:
        import docx
    except ImportError:
        logger.error("DOCX requires 'python-docx'. Run: uv pip install python-docx")
        return []

    logger.info(f"DOCX Loader: Parsing '{file_path}'")
    try:
        doc = docx.Document(file_path)
        documents = []
        current_section = []
        current_heading = "Document Start"

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            if para.style.name.startswith("Heading"):
                # Save the previous section
                if current_section:
                    documents.append(Document(
                        page_content="\n".join(current_section),
                        metadata={
                            "source": file_path,
                            "section": current_heading,
                            "file_type": "docx",
                        }
                    ))
                current_heading = text
                current_section = [f"## {text}"]
            else:
                current_section.append(text)

        # Don't forget the last section
        if current_section:
            documents.append(Document(
                page_content="\n".join(current_section),
                metadata={
                    "source": file_path,
                    "section": current_heading,
                    "file_type": "docx",
                }
            ))

        # Extract tables as separate Documents
        for i, table in enumerate(doc.tables):
            rows = []
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    rows.append(row_text)
            if rows:
                documents.append(Document(
                    page_content="\n".join(rows),
                    metadata={
                        "source": file_path,
                        "section": f"Table {i + 1}",
                        "file_type": "docx",
                        "element_type": "table",
                    }
                ))

        logger.info(f"DOCX Loader: Extracted {len(documents)} sections/tables")
        return documents

    except Exception as e:
        logger.error(f"DOCX Loader: Failed on '{file_path}': {e}")
        return []


# ──────────────────────────────────────────────────────────────────────────────
# XLSX — Excel Spreadsheets
# ──────────────────────────────────────────────────────────────────────────────
def parse_xlsx(file_path: str) -> List[Document]:
    """
    Parses an Excel .xlsx file.
    Returns one Document per sheet.
    """
    try:
        import openpyxl
    except ImportError:
        logger.error("XLSX requires 'openpyxl'. Run: uv pip install openpyxl")
        return []

    logger.info(f"XLSX Loader: Parsing '{file_path}'")
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        documents = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []

            for row in ws.iter_rows(values_only=True):
                row_values = [str(c) if c is not None else "" for c in row]
                if any(v.strip() for v in row_values):
                    rows.append(" | ".join(row_values))

            if rows:
                documents.append(Document(
                    page_content="\n".join(rows),
                    metadata={
                        "source": file_path,
                        "sheet": sheet_name,
                        "file_type": "xlsx",
                        "element_type": "table",
                    }
                ))

        wb.close()
        logger.info(f"XLSX Loader: Extracted {len(documents)} sheet(s)")
        return documents

    except Exception as e:
        logger.error(f"XLSX Loader: Failed on '{file_path}': {e}")
        return []


# ──────────────────────────────────────────────────────────────────────────────
# PPTX — PowerPoint Presentations
# ──────────────────────────────────────────────────────────────────────────────
def parse_pptx(file_path: str) -> List[Document]:
    """
    Parses a PowerPoint .pptx file.
    Returns one Document per slide.
    """
    try:
        from pptx import Presentation
    except ImportError:
        logger.error("PPTX requires 'python-pptx'. Run: uv pip install python-pptx")
        return []

    logger.info(f"PPTX Loader: Parsing '{file_path}'")
    try:
        prs = Presentation(file_path)
        documents = []
        total_slides = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, start=1):
            parts = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)

                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip()
                            for cell in row.cells
                            if cell.text.strip()
                        )
                        if row_text:
                            parts.append(row_text)

            if parts:
                slide_title = parts[0] if parts else f"Slide {slide_num}"
                documents.append(Document(
                    page_content="\n".join(parts),
                    metadata={
                        "source": file_path,
                        "slide": slide_num,
                        "total_slides": total_slides,
                        "slide_title": slide_title,
                        "file_type": "pptx",
                    }
                ))

        logger.info(f"PPTX Loader: Extracted {len(documents)} slide(s)")
        return documents

    except Exception as e:
        logger.error(f"PPTX Loader: Failed on '{file_path}': {e}")
        return []


# ──────────────────────────────────────────────────────────────────────────────
# CSV — Comma-Separated Values
# ──────────────────────────────────────────────────────────────────────────────
def parse_csv(file_path: str) -> List[Document]:
    """
    Parses a CSV file.
    Returns one Document per row, preserving column names as context.
    This gives better retrieval than treating the whole CSV as one chunk.
    """
    logger.info(f"CSV Loader: Parsing '{file_path}'")
    try:
        documents = []
        with open(file_path, "r", encoding="utf-8", errors="replace", newline="") as f:
            reader = csv.DictReader(f)
            headers = reader.fieldnames or []
            filename = os.path.basename(file_path)

            for row_num, row in enumerate(reader, start=1):
                row_text = "\n".join(
                    f"{k}: {v}" for k, v in row.items() if v and str(v).strip()
                )
                if row_text:
                    documents.append(Document(
                        page_content=row_text,
                        metadata={
                            "source": file_path,
                            "row": row_num,
                            "columns": headers,
                            "file_type": "csv",
                            "element_type": "table_row",
                        }
                    ))

        logger.info(f"CSV Loader: Extracted {len(documents)} rows")
        return documents

    except Exception as e:
        logger.error(f"CSV Loader: Failed on '{file_path}': {e}")
        return []
